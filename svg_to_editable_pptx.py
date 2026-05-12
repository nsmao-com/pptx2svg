#!/usr/bin/env python3
# input: SVG文件路径列表、背景渲染开关、首张页画布尺寸
# output: 单个可编辑PPTX文件，并提供SVG->PNG栅格化能力（含按页缩放、背景底色拆分、默认主题阴影清理、滤镜阴影拟合、圆角矩形原生化、字体名写入 latin/ea/cs 槽位、背景图去重缓存复用、转换总时长保护与噪点滤镜黑幕防护）
# pos: 2pptxsvg 中的 SVG->PPTX 核心转换模块
# 一旦我被更新，务必更新我的开头注释，以及所属的文件夹README。
"""
SVG -> Editable PPTX converter (MVP).

This script rebuilds SVG elements as native PowerPoint shapes so users can edit
individual elements in PPT.

Current coverage:
- Shapes: rect/circle/ellipse/path/polyline/polygon/line
- Text: text/tspan (basic)
- Image: image (local, data URI, http/https)
- Style: fill/stroke/opacity/font basic inheritance
- Transform: translate/scale/rotate/matrix/skew via accumulated matrix

Known limitations:
- Advanced filters/masks/blend/textPath are not supported
- Complex path fill rules and clipping are simplified
- Multi-stop gradients are reduced to 2 stops
"""

from __future__ import annotations

import argparse
import base64
import colorsys
import copy
import hashlib
import io
import math
import re
import sys
import time
import urllib.request
from dataclasses import dataclass, field
from pathlib import Path
from typing import Callable, Dict, Iterable, List, Optional, Tuple
from xml.etree import ElementTree as ET

import numpy as np
from PIL import Image, ImageColor
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Pt
from svgpathtools import Arc, CubicBezier, Line, QuadraticBezier, parse_path
from svgpathtools import parser as svg_parser

try:
    import cairosvg  # type: ignore

    HAS_CAIROSVG = True
except Exception:
    HAS_CAIROSVG = False

try:
    import skia  # type: ignore

    HAS_SKIA = True
except Exception:
    HAS_SKIA = False

EMU_PER_INCH = 914400
EMU_PER_PIXEL = EMU_PER_INCH / 96.0

XLINK_NS = "http://www.w3.org/1999/xlink"
SVG_NS = "http://www.w3.org/2000/svg"

ET.register_namespace("", SVG_NS)
ET.register_namespace("xlink", XLINK_NS)


@dataclass
class GradientStop:
    offset: float
    color: Tuple[int, int, int]
    opacity: float


@dataclass
class LinearGradient:
    id: str
    x1: float = 0.0
    y1: float = 0.0
    x2: float = 1.0
    y2: float = 0.0
    stops: List[GradientStop] = field(default_factory=list)


@dataclass
class SvgFilter:
    id: str
    std_deviation: float = 0.0
    has_source_over_blur: bool = False
    offset_dx: float = 0.0
    offset_dy: float = 0.0
    alpha_scale: float = 1.0
    shadow_rgb: Tuple[int, int, int] = (0, 0, 0)
    shadow_opacity: float = 1.0
    prefer_drop_shadow: bool = False


@dataclass
class RenderContext:
    slide: object
    vb_min_x: float
    vb_min_y: float
    vb_width: float
    vb_height: float
    target_width: float
    target_height: float
    scale_x: float
    scale_y: float
    offset_x: float
    offset_y: float
    stroke_scale: float
    svg_dir: Path
    gradients: Dict[str, LinearGradient]
    filters: Dict[str, SvgFilter]
    style_rules: Dict[str, Dict[str, Dict[str, str]]]
    skip_element_ids: set[int]


@dataclass
class BackgroundImageCache:
    png_by_key: Dict[str, bytes] = field(default_factory=dict)
    hits: int = 0
    renders: int = 0


DEFAULT_STYLE = {
    "display": "inline",
    "visibility": "visible",
    "color": "black",
    "fill": "black",
    "fill-opacity": "1",
    "stroke": "none",
    "stroke-opacity": "1",
    "stroke-width": "1",
    "font-size": "16",
    "font-family": "Calibri",
    "font-weight": "normal",
    "text-anchor": "start",
}

INHERITED_STYLE_KEYS = {
    "display",
    "visibility",
    "opacity",
    "color",
    "fill",
    "fill-opacity",
    "stroke",
    "stroke-opacity",
    "stroke-width",
    "stroke-dasharray",
    "font-size",
    "font-family",
    "font-weight",
    "text-anchor",
}


def emu(px: float) -> int:
    return int(round(px * EMU_PER_PIXEL))


def matrix_identity() -> np.ndarray:
    return np.eye(3)


def matrix_multiply(a: np.ndarray, b: np.ndarray) -> np.ndarray:
    return a @ b


def apply_matrix(m: np.ndarray, x: float, y: float) -> Tuple[float, float]:
    v = np.array([x, y, 1.0])
    out = m @ v
    return float(out[0]), float(out[1])


def parse_transform(transform: str | None) -> np.ndarray:
    if not transform:
        return matrix_identity()
    try:
        return np.array(svg_parser.parse_transform(transform), dtype=float)
    except Exception:
        return matrix_identity()


def parse_style_attr(style_attr: str | None) -> Dict[str, str]:
    if not style_attr:
        return {}
    out: Dict[str, str] = {}
    for chunk in style_attr.split(";"):
        if ":" not in chunk:
            continue
        k, v = chunk.split(":", 1)
        k = k.strip()
        v = v.strip()
        if k:
            out[k] = v
    return out


def parse_length(value: str | None, reference: float | None = None, default: float = 0.0) -> float:
    if value is None:
        return default
    text = str(value).strip()
    if not text:
        return default

    if text.endswith("%"):
        if reference is None:
            return default
        try:
            return float(text[:-1]) * reference / 100.0
        except ValueError:
            return default

    match = re.match(r"^([+-]?\d*\.?\d+(?:[eE][+-]?\d+)?)([a-zA-Z]*)$", text)
    if not match:
        return default

    num = float(match.group(1))
    unit = match.group(2).lower()

    if unit in ("", "px"):
        return num
    if unit == "pt":
        return num * 96.0 / 72.0
    if unit == "pc":
        return num * 16.0
    if unit == "in":
        return num * 96.0
    if unit == "cm":
        return num * 96.0 / 2.54
    if unit == "mm":
        return num * 96.0 / 25.4

    return num


def parse_float_list(value: str | None) -> List[float]:
    if not value:
        return []
    out: List[float] = []
    for part in re.split(r"[\s,]+", value.strip()):
        if not part:
            continue
        try:
            out.append(float(part))
        except ValueError:
            pass
    return out


def clamp01(v: float) -> float:
    return max(0.0, min(1.0, v))


def compute_canvas_mapping(
    vb_width: float,
    vb_height: float,
    target_width: float,
    target_height: float,
) -> Tuple[float, float, float, float, float]:
    vb_width = max(1.0, vb_width)
    vb_height = max(1.0, vb_height)
    target_width = max(1.0, target_width)
    target_height = max(1.0, target_height)

    scale = min(target_width / vb_width, target_height / vb_height)
    scaled_width = vb_width * scale
    scaled_height = vb_height * scale
    offset_x = (target_width - scaled_width) / 2.0
    offset_y = (target_height - scaled_height) / 2.0
    return scale, scale, offset_x, offset_y, scale


def parse_opacity_value(value: str | None, default: float = 1.0) -> float:
    if value is None:
        return clamp01(default)
    t = str(value).strip()
    if not t:
        return clamp01(default)
    try:
        if t.endswith("%"):
            return clamp01(float(t[:-1]) / 100.0)
        return clamp01(float(t))
    except ValueError:
        return clamp01(default)


def map_point_to_slide(ctx: RenderContext, x: float, y: float) -> Tuple[float, float]:
    slide_x = (x - ctx.vb_min_x) * ctx.scale_x + ctx.offset_x
    slide_y = (y - ctx.vb_min_y) * ctx.scale_y + ctx.offset_y
    return slide_x, slide_y


def _split_css_func_args(raw: str) -> List[str]:
    if "," in raw:
        return [x.strip() for x in raw.split(",") if x.strip()]
    compact = re.sub(r"\s*/\s*", " / ", raw.strip())
    return [x for x in compact.split() if x]


def parse_color(color_text: str | None) -> Optional[Tuple[int, int, int, float]]:
    if color_text is None:
        return None

    s = color_text.strip()
    if not s or s.lower() == "none":
        return None
    if s.lower().startswith("url("):
        return None
    if s.lower() == "transparent":
        return (0, 0, 0, 0.0)

    try:
        if s.lower().startswith("rgba") or s.lower().startswith("rgb"):
            m = re.match(r"rgba?\(([^)]+)\)", s, re.IGNORECASE)
            if not m:
                return None
            vals = _split_css_func_args(m.group(1))
            alpha = 1.0
            if "/" in vals:
                slash = vals.index("/")
                alpha = _parse_alpha(vals[slash + 1]) if slash + 1 < len(vals) else 1.0
                vals = vals[:slash]
            if len(vals) == 4:
                alpha = _parse_alpha(vals[3])
                vals = vals[:3]
            if len(vals) != 3:
                return None
            r = _parse_rgb_channel(vals[0])
            g = _parse_rgb_channel(vals[1])
            b = _parse_rgb_channel(vals[2])
            return (r, g, b, alpha)

        if s.lower().startswith("hsla") or s.lower().startswith("hsl"):
            m = re.match(r"hsla?\(([^)]+)\)", s, re.IGNORECASE)
            if not m:
                return None
            vals = _split_css_func_args(m.group(1))
            alpha = 1.0
            if "/" in vals:
                slash = vals.index("/")
                alpha = _parse_alpha(vals[slash + 1]) if slash + 1 < len(vals) else 1.0
                vals = vals[:slash]
            if len(vals) == 4:
                alpha = _parse_alpha(vals[3])
                vals = vals[:3]
            if len(vals) != 3:
                return None
            h = _parse_hue(vals[0])
            s1 = _parse_alpha(vals[1])
            l1 = _parse_alpha(vals[2])
            r, g, b = colorsys.hls_to_rgb(h, l1, s1)
            return (
                int(round(clamp01(r) * 255)),
                int(round(clamp01(g) * 255)),
                int(round(clamp01(b) * 255)),
                alpha,
            )

        rgb = ImageColor.getrgb(s)
        if isinstance(rgb, tuple) and len(rgb) == 3:
            return (int(rgb[0]), int(rgb[1]), int(rgb[2]), 1.0)
        if isinstance(rgb, tuple) and len(rgb) == 4:
            return (int(rgb[0]), int(rgb[1]), int(rgb[2]), clamp01(rgb[3] / 255.0))
    except Exception:
        return None

    return None


def _split_hex_alpha_color(color_text: str | None) -> Optional[Tuple[str, float]]:
    if color_text is None:
        return None
    s = color_text.strip()
    m8 = re.fullmatch(r"#([0-9a-fA-F]{8})", s)
    if m8:
        raw = m8.group(1)
        alpha = int(raw[6:8], 16) / 255.0
        return f"#{raw[:6]}", clamp01(alpha)

    m4 = re.fullmatch(r"#([0-9a-fA-F]{4})", s)
    if m4:
        raw = m4.group(1)
        rgb = "".join(ch * 2 for ch in raw[:3])
        alpha = int(raw[3] * 2, 16) / 255.0
        return f"#{rgb}", clamp01(alpha)

    return None


def _fmt_float(v: float) -> str:
    return f"{clamp01(v):.6f}".rstrip("0").rstrip(".") or "0"


def _hex_rgb(rgb: Tuple[int, int, int]) -> str:
    return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


def _parse_rgb_channel(v: str) -> int:
    t = v.strip()
    if t.endswith("%"):
        return max(0, min(255, int(round(float(t[:-1]) * 2.55))))
    return max(0, min(255, int(round(float(t)))))


def _parse_alpha(v: str) -> float:
    t = v.strip()
    if t.endswith("%"):
        return clamp01(float(t[:-1]) / 100.0)
    return clamp01(float(t))


def _parse_hue(v: str) -> float:
    t = v.strip().lower()
    if t.endswith("deg"):
        t = t[:-3].strip()
    try:
        deg = float(t)
    except ValueError:
        return 0.0
    return (deg % 360.0) / 360.0


def local_name(tag: str) -> str:
    if tag.startswith("{"):
        return tag.split("}", 1)[1]
    return tag


def parse_css_rules(root: ET.Element) -> Dict[str, Dict[str, Dict[str, str]]]:
    rules: Dict[str, Dict[str, Dict[str, str]]] = {
        "tag": {},
        "class": {},
        "id": {},
    }

    for elem in root.iter():
        if local_name(elem.tag) != "style":
            continue
        css_text = "".join(elem.itertext())
        if not css_text:
            continue
        css_text = re.sub(r"/\*.*?\*/", "", css_text, flags=re.DOTALL)
        css_text = re.sub(r"@import\s+[^;]+;", "", css_text, flags=re.IGNORECASE)
        css_text = re.sub(r"@charset\s+[^;]+;", "", css_text, flags=re.IGNORECASE)
        for selector_block, decl_block in re.findall(r"([^{}]+)\{([^{}]+)\}", css_text):
            declarations = parse_style_attr(decl_block)
            if not declarations:
                continue
            for selector in selector_block.split(","):
                s = selector.strip()
                if not s:
                    continue
                if re.fullmatch(r"[A-Za-z_][\w\-]*", s):
                    dst = rules["tag"].setdefault(s, {})
                    dst.update(declarations)
                elif re.fullmatch(r"\.[A-Za-z_][\w\-]*", s):
                    dst = rules["class"].setdefault(s[1:], {})
                    dst.update(declarations)
                elif re.fullmatch(r"#[A-Za-z_][\w\-]*", s):
                    dst = rules["id"].setdefault(s[1:], {})
                    dst.update(declarations)
    return rules


def _css_style_for_element(elem: ET.Element, style_rules: Dict[str, Dict[str, Dict[str, str]]]) -> Dict[str, str]:
    out: Dict[str, str] = {}
    tag = local_name(elem.tag)
    if tag in style_rules["tag"]:
        out.update(style_rules["tag"][tag])

    cls = elem.get("class", "")
    for name in cls.split():
        if name in style_rules["class"]:
            out.update(style_rules["class"][name])

    elem_id = elem.get("id")
    if elem_id and elem_id in style_rules["id"]:
        out.update(style_rules["id"][elem_id])
    return out


def merge_style(
    parent_style: Dict[str, str],
    elem: ET.Element,
    style_rules: Dict[str, Dict[str, Dict[str, str]]],
) -> Dict[str, str]:
    out = dict(parent_style)
    out.update(_css_style_for_element(elem, style_rules))
    out.update(parse_style_attr(elem.get("style")))
    for key in INHERITED_STYLE_KEYS:
        if key in elem.attrib:
            out[key] = elem.attrib[key]
    return out


def parse_viewbox(root: ET.Element) -> Tuple[float, float, float, float]:
    vb = root.get("viewBox")
    if vb:
        parts = parse_float_list(vb)
        if len(parts) >= 4 and parts[2] > 0 and parts[3] > 0:
            return parts[0], parts[1], parts[2], parts[3]

    width = parse_length(root.get("width"), default=1280)
    height = parse_length(root.get("height"), default=720)
    if width <= 0:
        width = 1280
    if height <= 0:
        height = 720
    return 0.0, 0.0, width, height


def parse_gradients(root: ET.Element, vb_width: float, vb_height: float) -> Dict[str, LinearGradient]:
    gradients: Dict[str, LinearGradient] = {}

    for elem in root.iter():
        tag = local_name(elem.tag)
        if tag not in {"linearGradient", "radialGradient"}:
            continue

        gid = elem.get("id")
        if not gid:
            continue

        grad = LinearGradient(id=gid)
        if tag == "linearGradient":
            grad.x1 = parse_length(elem.get("x1", "0%"), reference=vb_width, default=0.0)
            grad.y1 = parse_length(elem.get("y1", "0%"), reference=vb_height, default=0.0)
            grad.x2 = parse_length(elem.get("x2", "100%"), reference=vb_width, default=vb_width)
            grad.y2 = parse_length(elem.get("y2", "0%"), reference=vb_height, default=0.0)
        else:
            # Radial gradients are approximated using first/last stop colors.
            cx = parse_length(elem.get("cx", "50%"), reference=vb_width, default=vb_width / 2.0)
            cy = parse_length(elem.get("cy", "50%"), reference=vb_height, default=vb_height / 2.0)
            r = parse_length(
                elem.get("r", "50%"),
                reference=min(vb_width, vb_height),
                default=min(vb_width, vb_height) / 2.0,
            )
            grad.x1 = cx - r
            grad.y1 = cy
            grad.x2 = cx + r
            grad.y2 = cy

        stops: List[GradientStop] = []
        for stop in elem:
            if local_name(stop.tag) != "stop":
                continue
            offset_text = stop.get("offset", "0")
            if offset_text.endswith("%"):
                offset = clamp01(parse_length(offset_text, reference=100.0, default=0.0) / 100.0)
            else:
                try:
                    offset = clamp01(float(offset_text))
                except ValueError:
                    offset = 0.0

            stop_style = parse_style_attr(stop.get("style"))
            stop_color = stop.get("stop-color") or stop_style.get("stop-color") or "black"
            stop_opacity_text = stop.get("stop-opacity") or stop_style.get("stop-opacity") or "1"

            color = parse_color(stop_color)
            if not color:
                continue
            stop_opacity = parse_opacity_value(stop_opacity_text, default=1.0)

            stops.append(
                GradientStop(
                    offset=offset,
                    color=(color[0], color[1], color[2]),
                    opacity=clamp01(color[3] * stop_opacity),
                )
            )

        stops.sort(key=lambda s: s.offset)
        if stops:
            grad.stops = stops
            gradients[gid] = grad

    return gradients


def parse_filters(root: ET.Element) -> Dict[str, SvgFilter]:
    filters: Dict[str, SvgFilter] = {}

    for elem in root.iter():
        if local_name(elem.tag) != "filter":
            continue

        fid = (elem.get("id") or "").strip()
        if not fid:
            continue

        spec = SvgFilter(id=fid)
        for child in elem:
            tag = local_name(child.tag)
            if tag == "feGaussianBlur":
                nums = parse_float_list(child.get("stdDeviation"))
                if nums:
                    spec.std_deviation = max(spec.std_deviation, max(0.0, nums[0]))
                in_name = (child.get("in") or "").strip().lower()
                if in_name == "sourcealpha":
                    spec.prefer_drop_shadow = True
            elif tag == "feDropShadow":
                nums = parse_float_list(child.get("stdDeviation"))
                if nums:
                    spec.std_deviation = max(spec.std_deviation, max(0.0, nums[0]))
                spec.offset_dx = parse_length(child.get("dx"), default=spec.offset_dx)
                spec.offset_dy = parse_length(child.get("dy"), default=spec.offset_dy)
                color = parse_color(child.get("flood-color") or "black")
                if color:
                    spec.shadow_rgb = (color[0], color[1], color[2])
                    spec.shadow_opacity = clamp01(color[3])
                spec.shadow_opacity = clamp01(
                    spec.shadow_opacity
                    * parse_opacity_value(child.get("flood-opacity"), default=1.0)
                )
                spec.prefer_drop_shadow = True
            elif tag == "feOffset":
                spec.offset_dx = parse_length(child.get("dx"), default=spec.offset_dx)
                spec.offset_dy = parse_length(child.get("dy"), default=spec.offset_dy)
                spec.prefer_drop_shadow = True
            elif tag == "feComponentTransfer":
                for comp in child:
                    if local_name(comp.tag) != "feFuncA":
                        continue
                    mode = (comp.get("type") or "").strip().lower()
                    if mode == "linear":
                        slope = parse_length(comp.get("slope"), default=1.0)
                        intercept = parse_length(comp.get("intercept"), default=0.0)
                        spec.alpha_scale *= max(0.0, slope)
                        spec.shadow_opacity = clamp01(spec.shadow_opacity + max(0.0, intercept))
                    elif mode == "table":
                        table_vals = parse_float_list(comp.get("tableValues"))
                        if table_vals:
                            spec.alpha_scale *= max(0.0, table_vals[-1])
            elif tag == "feComposite":
                op = (child.get("operator") or "").strip().lower()
                in_a = (child.get("in") or "").strip().lower()
                in_b = (child.get("in2") or "").strip().lower()
                if op == "over" and ("sourcegraphic" in {in_a, in_b}):
                    spec.has_source_over_blur = True
                    spec.prefer_drop_shadow = True
            elif tag == "feMerge":
                for node in child:
                    if local_name(node.tag) != "feMergeNode":
                        continue
                    in_name = (node.get("in") or "").strip().lower()
                    if in_name == "sourcegraphic":
                        spec.has_source_over_blur = True
                        spec.prefer_drop_shadow = True

        if spec.std_deviation > 0:
            filters[fid] = spec

    return filters


RISKY_FILTER_PRIMITIVES_FOR_RASTER = {
    "feTurbulence",
    "feDisplacementMap",
    "feConvolveMatrix",
    "feMorphology",
    "feDiffuseLighting",
    "feSpecularLighting",
}


def _collect_risky_filter_ids(root: ET.Element) -> set[str]:
    risky: set[str] = set()
    for elem in root.iter():
        if local_name(elem.tag) != "filter":
            continue
        fid = (elem.get("id") or "").strip()
        if not fid:
            continue
        primitives = {local_name(child.tag) for child in elem if isinstance(child.tag, str)}
        if primitives & RISKY_FILTER_PRIMITIVES_FOR_RASTER:
            risky.add(fid)
    return risky


def _extract_gradient_id(fill_text: str | None) -> Optional[str]:
    if not fill_text:
        return None
    m = re.match(r"url\(#([^)]+)\)", fill_text.strip())
    if m:
        return m.group(1)
    return None


def _extract_filter_id(filter_text: str | None) -> Optional[str]:
    if not filter_text:
        return None
    m = re.match(r"url\(#([^)]+)\)", filter_text.strip())
    if m:
        return m.group(1)
    return None


def _has_explicit_fill_or_stroke(
    elem: ET.Element,
    own_style: Optional[Dict[str, str]] = None,
) -> bool:
    if elem.get("fill") is not None or elem.get("stroke") is not None:
        return True
    style_attr = parse_style_attr(elem.get("style"))
    if "fill" in style_attr or "stroke" in style_attr:
        return True
    if own_style and ("fill" in own_style or "stroke" in own_style):
        return True
    return False


def _is_rect_covering_viewbox(elem: ET.Element, vb_w: float, vb_h: float) -> bool:
    if local_name(elem.tag) != "rect":
        return False
    if elem.get("transform"):
        return False

    x = parse_length(elem.get("x"), reference=vb_w, default=0.0)
    y = parse_length(elem.get("y"), reference=vb_h, default=0.0)
    w = parse_length(elem.get("width"), reference=vb_w, default=0.0)
    h = parse_length(elem.get("height"), reference=vb_h, default=0.0)
    tol = max(1e-3, max(vb_w, vb_h) * 1e-3)
    if abs(x) > tol or abs(y) > tol:
        return False
    if w < vb_w - tol or h < vb_h - tol:
        return False
    return True


def _set_alpha_on_color(color_format, opacity: float) -> None:
    opacity = clamp01(opacity)
    if opacity >= 0.999:
        return
    xclr = getattr(color_format, "_xClr", None)
    if xclr is None:
        color_obj = getattr(color_format, "_color", None)
        xclr = getattr(color_obj, "_xClr", None)
    if xclr is None:
        return
    for child in list(xclr):
        if child.tag.endswith("}alpha"):
            xclr.remove(child)
    alpha = OxmlElement("a:alpha")
    alpha.set("val", str(int(round(opacity * 100000))))
    xclr.append(alpha)


def clear_shape_theme_effects(shape) -> None:
    element = getattr(shape, "_element", None)
    if element is None:
        return

    for child in list(element):
        if child.tag.endswith("}style"):
            element.remove(child)

    sp_pr = None
    for child in element:
        if child.tag.endswith("}spPr"):
            sp_pr = child
            break

    if sp_pr is None:
        return

    for child in list(sp_pr):
        if child.tag.endswith("}effectLst") or child.tag.endswith("}effectDag"):
            sp_pr.remove(child)

    effect_lst = OxmlElement("a:effectLst")
    insert_at = len(sp_pr)
    for idx, child in enumerate(sp_pr):
        local = child.tag.rsplit("}", 1)[-1]
        if local in {"scene3d", "sp3d", "extLst"}:
            insert_at = idx
            break
    sp_pr.insert(insert_at, effect_lst)


def _get_shape_effect_list(shape):
    element = getattr(shape, "_element", None)
    if element is None:
        return None

    sp_pr = None
    for child in element:
        if child.tag.endswith("}spPr"):
            sp_pr = child
            break
    if sp_pr is None:
        return None

    for child in sp_pr:
        if child.tag.endswith("}effectLst"):
            return child

    effect_lst = OxmlElement("a:effectLst")
    insert_at = len(sp_pr)
    for idx, child in enumerate(sp_pr):
        local = child.tag.rsplit("}", 1)[-1]
        if local in {"scene3d", "sp3d", "extLst"}:
            insert_at = idx
            break
    sp_pr.insert(insert_at, effect_lst)
    return effect_lst


def _pick_effect_color(style: Dict[str, str], gradients: Dict[str, LinearGradient]) -> Tuple[int, int, int, float]:
    fill_value = (style.get("fill") or "").strip()
    gid = _extract_gradient_id(fill_value)
    if gid and gid in gradients and gradients[gid].stops:
        stop = gradients[gid].stops[0]
        return (stop.color[0], stop.color[1], stop.color[2], stop.opacity)

    fill_color = parse_paint_value(fill_value, style)
    if fill_color:
        return fill_color

    stroke_color = parse_paint_value(style.get("stroke"), style)
    if stroke_color:
        return stroke_color

    return (255, 255, 255, 0.35)


def _remove_effect_nodes(effect_lst, locals_to_remove: set[str]) -> None:
    for child in list(effect_lst):
        local = child.tag.rsplit("}", 1)[-1]
        if local in locals_to_remove:
            effect_lst.remove(child)


def _append_color_with_alpha(parent, rgb: Tuple[int, int, int], alpha: float) -> None:
    srgb = OxmlElement("a:srgbClr")
    srgb.set("val", _hex_rgb(rgb))
    alpha = clamp01(alpha)
    if alpha < 0.999:
        alpha_elem = OxmlElement("a:alpha")
        alpha_elem.set("val", str(int(round(alpha * 100000))))
        srgb.append(alpha_elem)
    parent.append(srgb)


def _append_outer_shadow(
    effect_lst,
    rgb: Tuple[int, int, int],
    alpha: float,
    blur_px: float,
    dx_px: float,
    dy_px: float,
) -> None:
    shadow = OxmlElement("a:outerShdw")
    shadow.set("blurRad", str(int(round(max(0.25, blur_px) * EMU_PER_PIXEL))))
    dist_px = math.hypot(dx_px, dy_px)
    if dist_px > 1e-3:
        shadow.set("dist", str(int(round(dist_px * EMU_PER_PIXEL))))
        deg = math.degrees(math.atan2(dy_px, dx_px))
        if deg < 0:
            deg += 360.0
        shadow.set("dir", str(int(round(deg * 60000))))
    shadow.set("rotWithShape", "0")
    _append_color_with_alpha(shadow, rgb, alpha)
    effect_lst.append(shadow)


def apply_filter_effect_to_shape(
    shape,
    filter_id: Optional[str],
    ctx: RenderContext,
    style: Dict[str, str],
    opacity_multiplier: float,
) -> None:
    if not filter_id:
        return
    spec = ctx.filters.get(filter_id)
    if spec is None or spec.std_deviation <= 0:
        return

    effect_lst = _get_shape_effect_list(shape)
    if effect_lst is None:
        return

    _remove_effect_nodes(effect_lst, {"glow", "outerShdw"})

    # Prefer drop-shadow mapping when SVG filter carries blur + offset/source merge.
    if spec.prefer_drop_shadow or abs(spec.offset_dx) > 1e-3 or abs(spec.offset_dy) > 1e-3:
        alpha = clamp01(spec.shadow_opacity * spec.alpha_scale * opacity_multiplier)
        if alpha <= 1e-3:
            alpha = clamp01(0.18 * opacity_multiplier)
        blur_px = min(24.0, max(0.5, spec.std_deviation * 1.25 * ctx.stroke_scale))
        dx_px = spec.offset_dx * ctx.stroke_scale
        dy_px = spec.offset_dy * ctx.stroke_scale
        _append_outer_shadow(effect_lst, spec.shadow_rgb, alpha, blur_px, dx_px, dy_px)
        return

    color = _pick_effect_color(style, ctx.gradients)
    alpha = clamp01(color[3] * opacity_multiplier)
    if not spec.has_source_over_blur:
        alpha = clamp01(alpha * 0.55)

    # Fallback for non-shadow blur filters: keep a conservative glow.
    radius_px = min(20.0, max(0.5, spec.std_deviation * 1.25 * ctx.stroke_scale))
    glow = OxmlElement("a:glow")
    glow.set("rad", str(int(round(radius_px * EMU_PER_PIXEL))))
    _append_color_with_alpha(glow, (color[0], color[1], color[2]), alpha)
    effect_lst.append(glow)


def _set_run_fill_gradient(
    run,
    grad: LinearGradient,
    opacity: float,
) -> None:
    r_pr = run._r.get_or_add_rPr()
    for child in list(r_pr):
        if child.tag.endswith("}solidFill") or child.tag.endswith("}gradFill"):
            r_pr.remove(child)

    grad_fill = OxmlElement("a:gradFill")
    gs_lst = OxmlElement("a:gsLst")
    start = grad.stops[0]
    end = grad.stops[-1]

    for pos, stop in ((0, start), (100000, end)):
        gs = OxmlElement("a:gs")
        gs.set("pos", str(pos))
        srgb = OxmlElement("a:srgbClr")
        srgb.set("val", _hex_rgb(stop.color))
        stop_alpha = clamp01(stop.opacity * opacity)
        if stop_alpha < 0.999:
            alpha = OxmlElement("a:alpha")
            alpha.set("val", str(int(round(stop_alpha * 100000))))
            srgb.append(alpha)
        gs.append(srgb)
        gs_lst.append(gs)

    grad_fill.append(gs_lst)

    dx = grad.x2 - grad.x1
    dy = grad.y2 - grad.y1
    if abs(dx) > 1e-6 or abs(dy) > 1e-6:
        angle = math.degrees(math.atan2(dy, dx))
        # OOXML angle unit: 1/60000 degree.
        oo_angle = int(round(((450.0 - angle) % 360.0) * 60000))
        lin = OxmlElement("a:lin")
        lin.set("ang", str(oo_angle))
        lin.set("scaled", "1")
        grad_fill.append(lin)

    r_pr.append(grad_fill)


def _set_run_typefaces(run, font_name: str) -> None:
    font_name = (font_name or "").strip().strip("'").strip('"')
    if not font_name:
        return

    r_pr = run._r.get_or_add_rPr()
    for child in list(r_pr):
        if child.tag.endswith("}latin") or child.tag.endswith("}ea") or child.tag.endswith("}cs"):
            r_pr.remove(child)

    for tag_name in ("a:latin", "a:ea", "a:cs"):
        font_elem = OxmlElement(tag_name)
        font_elem.set("typeface", font_name)
        r_pr.append(font_elem)


def parse_paint_value(paint_text: str | None, style: Dict[str, str]) -> Optional[Tuple[int, int, int, float]]:
    if paint_text is None:
        return None
    text = paint_text.strip()
    if not text:
        return None
    if text.lower() == "currentcolor":
        text = style.get("color", "black")
    return parse_color(text)


def _fallback_color_from_paint_url(fill_value: str) -> Optional[str]:
    # Handles value like: url(#id) #fff
    m = re.match(r"url\(#([^)]+)\)\s*(.+)?", fill_value.strip())
    if not m:
        return None
    fallback = (m.group(2) or "").strip()
    return fallback or None


def apply_fill_style(
    shape,
    style: Dict[str, str],
    gradients: Dict[str, LinearGradient],
    opacity_multiplier: float = 1.0,
) -> None:
    fill_value = style.get("fill", "black").strip()
    fill_opacity = parse_opacity_value(style.get("fill-opacity"), default=1.0)
    total_opacity = clamp01(opacity_multiplier * fill_opacity)

    if fill_value.lower() == "none":
        shape.fill.background()
        return

    gid = _extract_gradient_id(fill_value)
    if gid and gid in gradients:
        grad = gradients[gid]
        shape.fill.gradient()

        start = grad.stops[0]
        end = grad.stops[-1]
        stops = shape.fill.gradient_stops

        stops[0].position = 0.0
        stops[0].color.rgb = RGBColor(*start.color)
        _set_alpha_on_color(stops[0].color, clamp01(start.opacity * total_opacity))

        stops[1].position = 1.0
        stops[1].color.rgb = RGBColor(*end.color)
        _set_alpha_on_color(stops[1].color, clamp01(end.opacity * total_opacity))

        dx = grad.x2 - grad.x1
        dy = grad.y2 - grad.y1
        if abs(dx) > 1e-6 or abs(dy) > 1e-6:
            angle = math.degrees(math.atan2(dy, dx))
            try:
                shape.fill.gradient_angle = float((450.0 - angle) % 360.0)
            except Exception:
                pass
        return

    color = parse_paint_value(fill_value, style)
    if not color:
        fallback = _fallback_color_from_paint_url(fill_value)
        color = parse_paint_value(fallback, style)
    if not color:
        shape.fill.background()
        return

    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(color[0], color[1], color[2])
    _set_alpha_on_color(shape.fill.fore_color, clamp01(color[3] * total_opacity))


def apply_line_style(
    shape,
    style: Dict[str, str],
    vb_ref: float,
    opacity_multiplier: float = 1.0,
    scale_multiplier: float = 1.0,
) -> None:
    stroke_value = style.get("stroke", "none").strip()
    stroke_opacity = parse_opacity_value(style.get("stroke-opacity"), default=1.0)
    total_opacity = clamp01(opacity_multiplier * stroke_opacity)

    if stroke_value.lower() == "none":
        shape.line.fill.background()
        return

    stroke_color = parse_paint_value(stroke_value, style)
    if not stroke_color:
        shape.line.fill.background()
        return

    stroke_width_px = parse_length(style.get("stroke-width", "1"), reference=vb_ref, default=1.0)
    if stroke_width_px < 0:
        stroke_width_px = 0
    stroke_width_px *= max(scale_multiplier, 1e-6)

    shape.line.width = Emu(emu(stroke_width_px))
    shape.line.color.rgb = RGBColor(stroke_color[0], stroke_color[1], stroke_color[2])
    _set_alpha_on_color(shape.line.color, clamp01(stroke_color[3] * total_opacity))

    dash = style.get("stroke-dasharray", "").strip()
    if dash and dash.lower() != "none":
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    else:
        shape.line.dash_style = MSO_LINE_DASH_STYLE.SOLID


def _sample_segment(segment, min_samples: int = 8) -> List[Tuple[float, float]]:
    if isinstance(segment, Line):
        return [
            (segment.start.real, segment.start.imag),
            (segment.end.real, segment.end.imag),
        ]

    if isinstance(segment, (CubicBezier, QuadraticBezier, Arc)):
        try:
            approx_len = max(1.0, segment.length(error=1e-3))
            n = max(min_samples, min(64, int(approx_len / 8.0)))
        except Exception:
            n = min_samples
        pts = []
        for i in range(n + 1):
            t = i / n
            p = segment.point(t)
            pts.append((p.real, p.imag))
        return pts

    return [
        (segment.start.real, segment.start.imag),
        (segment.end.real, segment.end.imag),
    ]


def path_to_contours(d: str) -> List[Tuple[List[Tuple[float, float]], bool]]:
    try:
        path = parse_path(d)
    except Exception:
        return []

    if len(path) == 0:
        return []

    contours: List[Tuple[List[Tuple[float, float]], bool]] = []
    current: List[Tuple[float, float]] = []
    start_pt: Optional[Tuple[float, float]] = None
    prev_end: Optional[complex] = None

    for seg in path:
        seg_start = complex(seg.start)
        seg_end = complex(seg.end)

        if prev_end is None or abs(seg_start - prev_end) > 1e-8:
            if current:
                contours.append((current, False))
            current = [(seg_start.real, seg_start.imag)]
            start_pt = (seg_start.real, seg_start.imag)

        sampled = _sample_segment(seg)
        if sampled:
            if current and _point_equal(current[-1], sampled[0]):
                current.extend(sampled[1:])
            else:
                current.extend(sampled)

        prev_end = seg_end

        if start_pt is not None and _point_equal((seg_end.real, seg_end.imag), start_pt):
            contours.append((current, True))
            current = []
            start_pt = None
            prev_end = None

    if current:
        contours.append((current, False))

    return contours


def _point_equal(p1: Tuple[float, float], p2: Tuple[float, float], tol: float = 1e-6) -> bool:
    return abs(p1[0] - p2[0]) <= tol and abs(p1[1] - p2[1]) <= tol


def transform_points(points: Iterable[Tuple[float, float]], m: np.ndarray, ctx: RenderContext) -> List[Tuple[float, float]]:
    out: List[Tuple[float, float]] = []
    for x, y in points:
        tx, ty = apply_matrix(m, x, y)
        out.append(map_point_to_slide(ctx, tx, ty))
    return out


def _dedupe_sequential_points(points: List[Tuple[float, float]], tol: float = 1e-4) -> List[Tuple[float, float]]:
    if not points:
        return points
    out = [points[0]]
    for p in points[1:]:
        if not _point_equal(out[-1], p, tol):
            out.append(p)
    return out


def render_freeform(
    ctx: RenderContext,
    contours: List[Tuple[List[Tuple[float, float]], bool]],
    style: Dict[str, str],
    transform_m: np.ndarray,
    opacity_multiplier: float,
    filter_id: Optional[str] = None,
) -> None:
    if not contours:
        return

    transformed_contours: List[Tuple[List[Tuple[float, float]], bool]] = []
    for pts, closed in contours:
        tpts = transform_points(pts, transform_m, ctx)
        tpts = _dedupe_sequential_points(tpts)
        if len(tpts) < 2:
            continue
        transformed_contours.append((tpts, closed))

    if not transformed_contours:
        return

    first_pts, _ = transformed_contours[0]
    fb = ctx.slide.shapes.build_freeform(
        start_x=first_pts[0][0],
        start_y=first_pts[0][1],
        scale=EMU_PER_PIXEL,
    )
    fb.add_line_segments(first_pts[1:], close=transformed_contours[0][1])

    for pts, closed in transformed_contours[1:]:
        fb.move_to(pts[0][0], pts[0][1])
        fb.add_line_segments(pts[1:], close=closed)

    shape = fb.convert_to_shape(Emu(0), Emu(0))
    apply_fill_style(shape, style, ctx.gradients, opacity_multiplier=opacity_multiplier)
    apply_line_style(
        shape,
        style,
        max(ctx.vb_width, ctx.vb_height),
        opacity_multiplier=opacity_multiplier,
        scale_multiplier=ctx.stroke_scale,
    )
    clear_shape_theme_effects(shape)
    apply_filter_effect_to_shape(shape, filter_id, ctx, style, opacity_multiplier)


def _render_rect_as_freeform(
    x: float,
    y: float,
    w: float,
    h: float,
    rx: float,
    ry: float,
    ctx: RenderContext,
    style: Dict[str, str],
    transform_m: np.ndarray,
    opacity_multiplier: float,
    filter_id: Optional[str],
) -> None:
    d = rect_to_path_d(x, y, w, h, rx, ry)
    render_freeform(
        ctx,
        path_to_contours(d),
        style,
        transform_m,
        opacity_multiplier,
        filter_id=filter_id,
    )


def _is_axis_aligned_scale_translate(m: np.ndarray, eps: float = 1e-6) -> bool:
    # [a c tx]
    # [b d ty]
    # [0 0 1 ]
    # Axis-aligned means no rotation/shear, i.e. b≈0 and c≈0.
    return abs(m[1, 0]) <= eps and abs(m[0, 1]) <= eps


def _render_rect_native_if_possible(
    x: float,
    y: float,
    w: float,
    h: float,
    rx: float,
    ry: float,
    ctx: RenderContext,
    style: Dict[str, str],
    transform_m: np.ndarray,
    opacity_multiplier: float,
    filter_id: Optional[str],
) -> bool:
    if not _is_axis_aligned_scale_translate(transform_m):
        return False

    a = float(transform_m[0, 0])
    d = float(transform_m[1, 1])
    if abs(a) <= 1e-8 or abs(d) <= 1e-8:
        return False

    # PPT rounded rectangle uses one radius parameter; if SVG has elliptical corners,
    # keep freeform fallback for fidelity.
    has_round = rx > 1e-6 or ry > 1e-6
    if has_round:
        if rx <= 1e-6:
            rx = ry
        if ry <= 1e-6:
            ry = rx
        if abs(rx - ry) > 1e-3:
            return False
        if abs(abs(a) - abs(d)) > 1e-6:
            return False

    p1x, p1y = apply_matrix(transform_m, x, y)
    p2x, p2y = apply_matrix(transform_m, x + w, y + h)
    left, top = map_point_to_slide(ctx, min(p1x, p2x), min(p1y, p2y))
    right, bottom = map_point_to_slide(ctx, max(p1x, p2x), max(p1y, p2y))
    out_w = max(1.0, right - left)
    out_h = max(1.0, bottom - top)

    shape_type = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if has_round else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    shape = ctx.slide.shapes.add_shape(
        shape_type,
        Emu(emu(left)),
        Emu(emu(top)),
        Emu(emu(out_w)),
        Emu(emu(out_h)),
    )

    if has_round:
        radius = max(rx, ry)
        radius = max(0.0, min(radius, w / 2.0, h / 2.0))
        min_side = max(1e-6, min(w, h))
        round_ratio = max(0.0, min(0.5, radius / min_side))
        try:
            shape.adjustments[0] = round_ratio
        except Exception:
            pass

    apply_fill_style(shape, style, ctx.gradients, opacity_multiplier=opacity_multiplier)
    apply_line_style(
        shape,
        style,
        max(ctx.vb_width, ctx.vb_height),
        opacity_multiplier=opacity_multiplier,
        scale_multiplier=ctx.stroke_scale,
    )
    clear_shape_theme_effects(shape)
    apply_filter_effect_to_shape(shape, filter_id, ctx, style, opacity_multiplier)
    return True


def render_rect(
    elem: ET.Element,
    ctx: RenderContext,
    style: Dict[str, str],
    m: np.ndarray,
    opacity_multiplier: float,
    filter_id: Optional[str] = None,
) -> None:
    x = parse_length(elem.get("x"), reference=ctx.vb_width, default=0.0)
    y = parse_length(elem.get("y"), reference=ctx.vb_height, default=0.0)
    w = parse_length(elem.get("width"), reference=ctx.vb_width, default=0.0)
    h = parse_length(elem.get("height"), reference=ctx.vb_height, default=0.0)
    if w <= 0 or h <= 0:
        return

    rx = parse_length(elem.get("rx"), reference=w, default=0.0)
    ry = parse_length(elem.get("ry"), reference=h, default=0.0)

    if _render_rect_native_if_possible(
        x,
        y,
        w,
        h,
        rx,
        ry,
        ctx,
        style,
        m,
        opacity_multiplier,
        filter_id,
    ):
        return

    _render_rect_as_freeform(
        x,
        y,
        w,
        h,
        rx,
        ry,
        ctx,
        style,
        m,
        opacity_multiplier,
        filter_id,
    )


def rect_to_path_d(x: float, y: float, width: float, height: float, rx: float, ry: float) -> str:
    rx = max(0.0, min(rx, width / 2.0))
    ry = max(0.0, min(ry, height / 2.0))

    if rx <= 0 and ry <= 0:
        return f"M{x},{y} L{x + width},{y} L{x + width},{y + height} L{x},{y + height} Z"

    if rx <= 0:
        rx = ry
    if ry <= 0:
        ry = rx

    x1 = x + rx
    x2 = x + width - rx
    y1 = y + ry
    y2 = y + height - ry

    return (
        f"M{x1},{y} "
        f"L{x2},{y} "
        f"A{rx},{ry} 0 0 1 {x + width},{y1} "
        f"L{x + width},{y2} "
        f"A{rx},{ry} 0 0 1 {x2},{y + height} "
        f"L{x1},{y + height} "
        f"A{rx},{ry} 0 0 1 {x},{y2} "
        f"L{x},{y1} "
        f"A{rx},{ry} 0 0 1 {x1},{y} Z"
    )


def points_attr_to_pairs(points_text: str | None) -> List[Tuple[float, float]]:
    nums = parse_float_list(points_text)
    if len(nums) < 2:
        return []
    pairs: List[Tuple[float, float]] = []
    for i in range(0, len(nums) - 1, 2):
        pairs.append((nums[i], nums[i + 1]))
    return pairs


def _matrix_rotation_deg(m: np.ndarray) -> float:
    a = m[0, 0]
    b = m[1, 0]
    return math.degrees(math.atan2(b, a))


def _matrix_has_shear(m: np.ndarray, eps: float = 1e-4) -> bool:
    a, c = m[0, 0], m[0, 1]
    b, d = m[1, 0], m[1, 1]
    dot = a * c + b * d
    return abs(dot) > eps


def render_text(
    elem: ET.Element,
    ctx: RenderContext,
    style: Dict[str, str],
    m: np.ndarray,
    opacity_multiplier: float,
    filter_id: Optional[str] = None,
) -> None:
    x_vals = parse_float_list(elem.get("x"))
    y_vals = parse_float_list(elem.get("y"))
    x = x_vals[0] if x_vals else parse_length(elem.get("x"), reference=ctx.vb_width, default=0.0)
    y = y_vals[0] if y_vals else parse_length(elem.get("y"), reference=ctx.vb_height, default=0.0)

    line_runs: List[List[Tuple[str, Dict[str, str]]]] = []
    if any(local_name(c.tag) == "tspan" for c in elem):
        current_runs: List[Tuple[str, Dict[str, str]]] = []

        def flush_current() -> None:
            nonlocal current_runs
            if any(seg.strip() for seg, _ in current_runs):
                line_runs.append(current_runs)
            current_runs = []

        lead_text = (elem.text or "").strip()
        if lead_text:
            current_runs.append((lead_text, style))

        for child in elem:
            if local_name(child.tag) != "tspan":
                tail = (child.tail or "").strip()
                if tail:
                    current_runs.append((tail, style))
                continue

            child_style = merge_style(style, child, ctx.style_rules)
            child_text = "".join(child.itertext()).strip()
            if not child_text:
                tail = (child.tail or "").strip()
                if tail:
                    current_runs.append((tail, style))
                continue

            # tspan with explicit x/y/dy typically indicates a new line chunk.
            if child.get("x") is not None or child.get("y") is not None or child.get("dy") is not None:
                flush_current()

            dx = parse_length(child.get("dx"), reference=ctx.vb_width, default=0.0)
            if dx > 0:
                base_size = parse_length(child_style.get("font-size", "16"), reference=ctx.vb_height, default=16.0)
                est_char_w = max(1.0, base_size * 0.5)
                spaces = min(24, max(0, int(round(dx / est_char_w))))
                if spaces > 0:
                    child_text = (" " * spaces) + child_text

            current_runs.append((child_text, child_style))

            tail = (child.tail or "").strip()
            if tail:
                current_runs.append((tail, style))

        flush_current()
    else:
        t = "".join(elem.itertext()).strip()
        if t:
            line_runs.append([(t, style)])

    if not line_runs:
        return

    font_size_px = parse_length(style.get("font-size", "16"), reference=ctx.vb_height, default=16.0)
    anchor = style.get("text-anchor", "start").strip().lower()

    line_texts = ["".join(seg for seg, _ in runs) for runs in line_runs]
    max_len = max(len(line) for line in line_texts)
    text_w = max(8.0, max_len * font_size_px * 0.6)
    text_h = max(font_size_px * 1.4, len(line_runs) * font_size_px * 1.35)

    if anchor == "middle":
        x -= text_w / 2.0
    elif anchor == "end":
        x -= text_w

    tx, ty = apply_matrix(m, x, y - font_size_px)
    tx, ty = map_point_to_slide(ctx, tx, ty)
    out_w = max(8.0, text_w * ctx.scale_x)
    out_h = max(8.0, text_h * ctx.scale_y)

    tb = ctx.slide.shapes.add_textbox(Emu(emu(tx)), Emu(emu(ty)), Emu(emu(out_w)), Emu(emu(out_h)))
    clear_shape_theme_effects(tb)
    tf = tb.text_frame
    tf.clear()

    for idx, runs in enumerate(line_runs):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        for seg_text, seg_style in runs:
            run = p.add_run()
            run.text = seg_text

            seg_font_size_px = parse_length(seg_style.get("font-size", "16"), reference=ctx.vb_height, default=font_size_px)
            seg_font_name = seg_style.get("font-family", "Calibri").split(",")[0].strip().strip("'").strip('"')
            seg_font_weight = seg_style.get("font-weight", "normal").strip().lower()

            run.font.size = Pt(seg_font_size_px * ctx.stroke_scale * 72.0 / 96.0)
            run.font.name = seg_font_name
            _set_run_typefaces(run, seg_font_name)
            run.font.bold = seg_font_weight in {"bold", "700", "800", "900"}

            fill_value = seg_style.get("fill", "black")
            color = parse_paint_value(fill_value, seg_style)
            if not color:
                fallback = _fallback_color_from_paint_url(fill_value)
                color = parse_paint_value(fallback, seg_style)
            fill_opacity = parse_opacity_value(seg_style.get("fill-opacity"), default=1.0)
            gradient_id = _extract_gradient_id(fill_value)
            text_gradient = ctx.gradients.get(gradient_id) if gradient_id else None

            if text_gradient and text_gradient.stops:
                _set_run_fill_gradient(
                    run,
                    text_gradient,
                    clamp01(opacity_multiplier * fill_opacity),
                )
            elif color:
                run.font.color.rgb = RGBColor(color[0], color[1], color[2])
                _set_alpha_on_color(
                    run.font.color,
                    clamp01(color[3] * opacity_multiplier * fill_opacity),
                )

        if anchor == "middle":
            p.alignment = PP_ALIGN.CENTER
        elif anchor == "end":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT

    if not _matrix_has_shear(m):
        tb.rotation = _matrix_rotation_deg(m)

    apply_filter_effect_to_shape(tb, filter_id, ctx, style, opacity_multiplier)


def resolve_image_source(href: str, ctx: RenderContext) -> Optional[str | io.BytesIO]:
    href = href.strip()
    if not href:
        return None

    if href.startswith("data:"):
        m = re.match(r"data:[^;]+;base64,(.*)$", href, re.IGNORECASE | re.DOTALL)
        if not m:
            return None
        try:
            blob = base64.b64decode(m.group(1), validate=False)
            return io.BytesIO(blob)
        except Exception:
            return None

    if href.startswith("http://") or href.startswith("https://"):
        try:
            req = urllib.request.Request(href, headers={"User-Agent": "svg2pptx-editable"})
            with urllib.request.urlopen(req, timeout=20) as resp:
                blob = resp.read()
            return io.BytesIO(blob)
        except Exception:
            return None

    p = Path(href)
    if not p.is_absolute():
        p = ctx.svg_dir / p
    if p.exists():
        return str(p)
    return None


def render_image(elem: ET.Element, ctx: RenderContext, m: np.ndarray) -> None:
    href = elem.get(f"{{{XLINK_NS}}}href") or elem.get("href")
    if not href:
        return

    x = parse_length(elem.get("x"), reference=ctx.vb_width, default=0.0)
    y = parse_length(elem.get("y"), reference=ctx.vb_height, default=0.0)
    w = parse_length(elem.get("width"), reference=ctx.vb_width, default=0.0)
    h = parse_length(elem.get("height"), reference=ctx.vb_height, default=0.0)
    if w <= 0 or h <= 0:
        return

    corners = [(x, y), (x + w, y), (x + w, y + h), (x, y + h)]
    tc = transform_points(corners, m, ctx)
    xs = [p[0] for p in tc]
    ys = [p[1] for p in tc]
    left = min(xs)
    top = min(ys)
    out_w = max(1.0, max(xs) - min(xs))
    out_h = max(1.0, max(ys) - min(ys))

    src = resolve_image_source(href, ctx)
    if src is None:
        return

    try:
        pic = ctx.slide.shapes.add_picture(
            src,
            Emu(emu(left)),
            Emu(emu(top)),
            Emu(emu(out_w)),
            Emu(emu(out_h)),
        )
        if not _matrix_has_shear(m):
            pic.rotation = _matrix_rotation_deg(m)
    except Exception:
        return


def render_line(
    elem: ET.Element,
    ctx: RenderContext,
    style: Dict[str, str],
    m: np.ndarray,
    opacity_multiplier: float,
    filter_id: Optional[str] = None,
) -> None:
    x1 = parse_length(elem.get("x1"), reference=ctx.vb_width, default=0.0)
    y1 = parse_length(elem.get("y1"), reference=ctx.vb_height, default=0.0)
    x2 = parse_length(elem.get("x2"), reference=ctx.vb_width, default=0.0)
    y2 = parse_length(elem.get("y2"), reference=ctx.vb_height, default=0.0)

    p1 = transform_points([(x1, y1)], m, ctx)[0]
    p2 = transform_points([(x2, y2)], m, ctx)[0]

    line = ctx.slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Emu(emu(p1[0])),
        Emu(emu(p1[1])),
        Emu(emu(p2[0])),
        Emu(emu(p2[1])),
    )
    clear_shape_theme_effects(line)

    style2 = dict(style)
    style2["fill"] = "none"
    apply_line_style(
        line,
        style2,
        max(ctx.vb_width, ctx.vb_height),
        opacity_multiplier=opacity_multiplier,
        scale_multiplier=ctx.stroke_scale,
    )
    apply_filter_effect_to_shape(line, filter_id, ctx, style, opacity_multiplier)


def _should_skip_implicit_filtered_canvas_rect(
    elem: ET.Element,
    tag: str,
    own_style: Dict[str, str],
    filter_id: Optional[str],
    ctx: RenderContext,
) -> bool:
    if tag != "rect" or not filter_id:
        return False
    # Only keep the defensive skip for unsupported filters. Gaussian blur filters
    # are represented in ctx.filters and can be approximated by PPT effects.
    if filter_id in ctx.filters:
        return False
    if _has_explicit_fill_or_stroke(elem, own_style):
        return False
    return _is_rect_covering_viewbox(elem, ctx.vb_width, ctx.vb_height)


def render_element(
    elem: ET.Element,
    ctx: RenderContext,
    parent_style: Dict[str, str],
    parent_m: np.ndarray,
    parent_opacity: float,
) -> None:
    if id(elem) in ctx.skip_element_ids:
        return

    tag = local_name(elem.tag)
    own_style = merge_style({}, elem, ctx.style_rules)
    style = merge_style(parent_style, elem, ctx.style_rules)

    if style.get("display", "inline") == "none":
        return
    if style.get("visibility", "visible") == "hidden":
        return

    local_opacity = parse_opacity_value(own_style.get("opacity"), default=1.0)
    effective_opacity = clamp01(parent_opacity * local_opacity)
    filter_id = _extract_filter_id(elem.get("filter") or own_style.get("filter"))

    current_m = matrix_multiply(parent_m, parse_transform(elem.get("transform")))

    if _should_skip_implicit_filtered_canvas_rect(elem, tag, own_style, filter_id, ctx):
        return

    if tag in {"svg", "g", "a"}:
        for child in elem:
            if isinstance(child.tag, str):
                render_element(child, ctx, style, current_m, effective_opacity)
        return

    if tag in {"defs", "style", "script", "metadata", "title", "desc", "clipPath", "mask", "filter"}:
        return

    if tag == "rect":
        render_rect(elem, ctx, style, current_m, effective_opacity, filter_id=filter_id)
        return

    if tag == "circle":
        cx = parse_length(elem.get("cx"), reference=ctx.vb_width, default=0.0)
        cy = parse_length(elem.get("cy"), reference=ctx.vb_height, default=0.0)
        r = parse_length(elem.get("r"), reference=min(ctx.vb_width, ctx.vb_height), default=0.0)
        if r <= 0:
            return
        d = f"M {cx - r},{cy} A {r},{r} 0 1 0 {cx + r},{cy} A {r},{r} 0 1 0 {cx - r},{cy} Z"
        render_freeform(
            ctx,
            path_to_contours(d),
            style,
            current_m,
            effective_opacity,
            filter_id=filter_id,
        )
        return

    if tag == "ellipse":
        cx = parse_length(elem.get("cx"), reference=ctx.vb_width, default=0.0)
        cy = parse_length(elem.get("cy"), reference=ctx.vb_height, default=0.0)
        rx = parse_length(elem.get("rx"), reference=ctx.vb_width, default=0.0)
        ry = parse_length(elem.get("ry"), reference=ctx.vb_height, default=0.0)
        if rx <= 0 or ry <= 0:
            return
        d = f"M {cx - rx},{cy} A {rx},{ry} 0 1 0 {cx + rx},{cy} A {rx},{ry} 0 1 0 {cx - rx},{cy} Z"
        render_freeform(
            ctx,
            path_to_contours(d),
            style,
            current_m,
            effective_opacity,
            filter_id=filter_id,
        )
        return

    if tag == "line":
        render_line(elem, ctx, style, current_m, effective_opacity, filter_id=filter_id)
        return

    if tag == "polyline":
        pts = points_attr_to_pairs(elem.get("points"))
        if len(pts) < 2:
            return
        render_freeform(ctx, [(pts, False)], style, current_m, effective_opacity, filter_id=filter_id)
        return

    if tag == "polygon":
        pts = points_attr_to_pairs(elem.get("points"))
        if len(pts) < 3:
            return
        render_freeform(ctx, [(pts, True)], style, current_m, effective_opacity, filter_id=filter_id)
        return

    if tag == "path":
        d = elem.get("d")
        if not d:
            return
        render_freeform(
            ctx,
            path_to_contours(d),
            style,
            current_m,
            effective_opacity,
            filter_id=filter_id,
        )
        return

    if tag == "text":
        render_text(elem, ctx, style, current_m, effective_opacity, filter_id=filter_id)
        return

    if tag == "image":
        render_image(elem, ctx, current_m)
        return

    for child in elem:
        if isinstance(child.tag, str):
            render_element(child, ctx, style, current_m, effective_opacity)


def find_background_group(root: ET.Element) -> Optional[ET.Element]:
    preferred_ids = {"bg", "background", "bg_layer", "background_layer"}
    for child in root:
        if local_name(child.tag) != "g":
            continue
        gid = (child.get("id") or "").strip().lower()
        if gid in preferred_ids:
            return child
    for child in root:
        if local_name(child.tag) != "g":
            continue
        gid = (child.get("id") or "").strip().lower()
        if "bg" in gid or "background" in gid:
            return child
    return None


def _render_svg_bytes_to_png(svg_bytes: bytes, out_w: float, out_h: float) -> Optional[bytes]:
    out_w = max(1, int(round(out_w)))
    out_h = max(1, int(round(out_h)))
    svg_bytes = _patch_svg_font_stack_for_raster(svg_bytes)

    if HAS_CAIROSVG:
        try:
            return cairosvg.svg2png(
                bytestring=svg_bytes,
                output_width=out_w,
                output_height=out_h,
            )
        except Exception:
            pass

    if HAS_SKIA:
        try:
            stream = skia.MemoryStream(svg_bytes)
            dom = skia.SVGDOM.MakeFromStream(stream)
            if dom is None:
                return None
            dom.setContainerSize(skia.Size.Make(out_w, out_h))
            surface = skia.Surface(out_w, out_h)
            canvas = surface.getCanvas()
            canvas.clear(skia.ColorTRANSPARENT)
            dom.render(canvas)
            img = surface.makeImageSnapshot()
            data = img.encodeToData()
            if data is None:
                return None
            return bytes(data)
        except Exception:
            return None

    return None


def _normalize_hex_alpha_colors_for_raster(svg_text: str) -> str:
    try:
        root = ET.fromstring(svg_text)
    except Exception:
        return svg_text

    paint_to_opacity = {
        "fill": "fill-opacity",
        "stroke": "stroke-opacity",
        "stop-color": "stop-opacity",
        "flood-color": "flood-opacity",
    }
    changed = False

    for elem in root.iter():
        for paint_attr, opacity_attr in paint_to_opacity.items():
            parsed = _split_hex_alpha_color(elem.get(paint_attr))
            if parsed is None:
                continue
            base_color, alpha = parsed
            elem.set(paint_attr, base_color)
            current_opacity = parse_opacity_value(elem.get(opacity_attr), default=1.0)
            elem.set(opacity_attr, _fmt_float(current_opacity * alpha))
            changed = True

        style_attr = elem.get("style")
        if not style_attr:
            continue
        style_map = parse_style_attr(style_attr)
        style_changed = False
        for paint_attr, opacity_attr in paint_to_opacity.items():
            parsed = _split_hex_alpha_color(style_map.get(paint_attr))
            if parsed is None:
                continue
            base_color, alpha = parsed
            style_map[paint_attr] = base_color
            current_opacity = parse_opacity_value(style_map.get(opacity_attr), default=1.0)
            style_map[opacity_attr] = _fmt_float(current_opacity * alpha)
            style_changed = True
        if style_changed:
            elem.set("style", "; ".join(f"{k}: {v}" for k, v in style_map.items()))
            changed = True

        if local_name(elem.tag) == "style" and elem.text:
            css_text = elem.text

            def _repl_hex8(match: re.Match[str]) -> str:
                raw = match.group(1)
                r = int(raw[0:2], 16)
                g = int(raw[2:4], 16)
                b = int(raw[4:6], 16)
                a = int(raw[6:8], 16) / 255.0
                return f"rgba({r},{g},{b},{_fmt_float(a)})"

            def _repl_hex4(match: re.Match[str]) -> str:
                raw = match.group(1)
                r = int(raw[0] * 2, 16)
                g = int(raw[1] * 2, 16)
                b = int(raw[2] * 2, 16)
                a = int(raw[3] * 2, 16) / 255.0
                return f"rgba({r},{g},{b},{_fmt_float(a)})"

            css_text2 = re.sub(r"#([0-9a-fA-F]{8})\b", _repl_hex8, css_text)
            css_text2 = re.sub(r"#([0-9a-fA-F]{4})\b", _repl_hex4, css_text2)
            if css_text2 != css_text:
                elem.text = css_text2
                changed = True

    if not changed:
        return svg_text

    return ET.tostring(root, encoding="unicode")


def _strip_problematic_overlay_rects_for_raster(svg_text: str) -> str:
    try:
        root = ET.fromstring(svg_text)
    except Exception:
        return svg_text

    risky_filters = _collect_risky_filter_ids(root)
    if not risky_filters:
        return svg_text

    _, _, vb_w, vb_h = parse_viewbox(root)
    style_rules = parse_css_rules(root)
    parent_by_child_id = {id(child): parent for parent in root.iter() for child in parent}
    changed = False

    for elem in list(root.iter()):
        if local_name(elem.tag) != "rect":
            continue
        own_style = merge_style({}, elem, style_rules)
        if _has_explicit_fill_or_stroke(elem, own_style):
            continue
        filter_id = _extract_filter_id(elem.get("filter") or own_style.get("filter"))
        if not filter_id or filter_id not in risky_filters:
            continue
        if not _is_rect_covering_viewbox(elem, vb_w, vb_h):
            continue
        parent = parent_by_child_id.get(id(elem))
        if parent is None:
            continue
        parent.remove(elem)
        changed = True

    if not changed:
        return svg_text

    return ET.tostring(root, encoding="unicode")


def _patch_svg_font_stack_for_raster(svg_bytes: bytes) -> bytes:
    try:
        svg_text = svg_bytes.decode("utf-8")
    except UnicodeDecodeError:
        svg_text = svg_bytes.decode("utf-8", errors="ignore")

    # Normalize 8-digit hex colors into color + explicit opacity for raster backends.
    svg_text = _normalize_hex_alpha_colors_for_raster(svg_text)
    # Drop full-canvas implicit black rects behind complex filters (e.g. feTurbulence)
    # to prevent backend-specific "black overlay" artifacts.
    svg_text = _strip_problematic_overlay_rects_for_raster(svg_text)

    # Remove remote @import to avoid render drift/timeouts in isolated servers.
    svg_text = re.sub(
        r"@import\s+url\([^)]+\)\s*;",
        "",
        svg_text,
        flags=re.IGNORECASE,
    )

    # Ensure Chinese-friendly fallback even when webfont declarations are unavailable.
    fallback_stack = (
        "'HarmonyOS Sans SC', 'Noto Sans CJK SC', 'Noto Sans SC', "
        "'Source Han Sans SC', 'Microsoft YaHei', 'PingFang SC', "
        "'WenQuanYi Micro Hei', sans-serif"
    )
    svg_text = re.sub(
        r'font-family\s*:\s*([\'"])?HarmonyOS\s+Sans\s+SC\1\s*;',
        f"font-family: {fallback_stack};",
        svg_text,
        flags=re.IGNORECASE,
    )
    svg_text = re.sub(
        r'font-family=(["\'])HarmonyOS\s+Sans\s+SC\1',
        f'font-family="{fallback_stack}"',
        svg_text,
        flags=re.IGNORECASE,
    )
    svg_text = _inject_font_style_for_known_classes(svg_text, fallback_stack)

    inject_css = f"text, tspan {{ font-family: {fallback_stack}; }}"
    style_block = f"<style><![CDATA[{inject_css}]]></style>"

    # Prefer appending to existing <defs>, otherwise inject right after <svg ...>.
    defs_match = re.search(r"<defs[^>]*>", svg_text, flags=re.IGNORECASE)
    if defs_match:
        idx = defs_match.end()
        svg_text = svg_text[:idx] + style_block + svg_text[idx:]
    else:
        svg_match = re.search(r"<svg[^>]*>", svg_text, flags=re.IGNORECASE)
        if svg_match:
            idx = svg_match.end()
            svg_text = svg_text[:idx] + style_block + svg_text[idx:]

    return svg_text.encode("utf-8")


def _inject_font_style_for_known_classes(svg_text: str, fallback_stack: str) -> str:
    class_style_map = {
        "ppt-font-cn-heading": f"font-family: {fallback_stack}; font-weight: 700;",
        "ppt-font-cn-body": f"font-family: {fallback_stack}; font-weight: 400;",
    }
    tag_pattern = re.compile(r"<(text|tspan)\b([^>]*)>", flags=re.IGNORECASE)

    def _apply(match: re.Match[str]) -> str:
        tag = match.group(1)
        attrs = match.group(2)

        class_match = re.search(
            r'\bclass=(["\'])(.*?)\1',
            attrs,
            flags=re.IGNORECASE | re.DOTALL,
        )
        if not class_match:
            return match.group(0)

        classes = {
            token.strip()
            for token in re.split(r"\s+", class_match.group(2).strip())
            if token.strip()
        }
        extra_style_parts = [
            css for cls, css in class_style_map.items() if cls in classes
        ]
        if not extra_style_parts:
            return match.group(0)

        extra_style = " ".join(extra_style_parts)
        style_match = re.search(
            r'\bstyle=(["\'])(.*?)\1',
            attrs,
            flags=re.IGNORECASE | re.DOTALL,
        )
        if style_match:
            quote = style_match.group(1)
            current_style = style_match.group(2).strip()
            sep = "; " if current_style and not current_style.endswith(";") else " "
            merged_style = f"{current_style}{sep}{extra_style}".strip()
            new_style_attr = f"style={quote}{merged_style}{quote}"
            attrs = attrs[: style_match.start()] + new_style_attr + attrs[style_match.end() :]
        else:
            attrs = f'{attrs} style="{extra_style}"'

        return f"<{tag}{attrs}>"

    return tag_pattern.sub(_apply, svg_text)


def parse_svg_canvas_size(svg_content: str) -> Tuple[float, float]:
    root = ET.fromstring(svg_content)
    _, _, vb_w, vb_h = parse_viewbox(root)
    return vb_w, vb_h


def rasterize_svg_to_png(svg_content: str, output_width: int, output_height: int) -> bytes:
    if not isinstance(svg_content, str) or not svg_content.strip():
        raise ValueError("svg content is empty")
    if output_width <= 0 or output_height <= 0:
        raise ValueError("output width/height must be positive")

    png_bytes = _render_svg_bytes_to_png(
        svg_content.encode("utf-8"),
        float(output_width),
        float(output_height),
    )
    if png_bytes is None:
        if not (HAS_CAIROSVG or HAS_SKIA):
            raise RuntimeError("no SVG raster backend available")
        raise RuntimeError("svg rasterization failed")
    return png_bytes


def _clone_without_elements(elem: ET.Element, omitted_ids: set[int]) -> Optional[ET.Element]:
    if id(elem) in omitted_ids:
        return None

    clone = ET.Element(elem.tag, dict(elem.attrib))
    clone.text = elem.text
    clone.tail = elem.tail
    for child in elem:
        cloned_child = _clone_without_elements(child, omitted_ids)
        if cloned_child is not None:
            clone.append(cloned_child)
    return clone


def _build_background_svg_bytes(
    root: ET.Element,
    bg_group: ET.Element,
    vb_w: float,
    vb_h: float,
    omitted_ids: Optional[set[int]] = None,
) -> bytes:
    svg = ET.Element(
        f"{{{SVG_NS}}}svg",
        {
            "viewBox": root.get("viewBox") or f"0 0 {vb_w} {vb_h}",
            "width": root.get("width") or str(vb_w),
            "height": root.get("height") or str(vb_h),
        },
    )

    omitted_ids = omitted_ids or set()

    # Preserve global defs/style for gradient/pattern/filter dependencies.
    for child in root:
        lname = local_name(child.tag)
        if lname in {"defs", "style"}:
            svg.append(copy.deepcopy(child))
            continue
        if child is bg_group:
            cloned_bg = _clone_without_elements(bg_group, omitted_ids)
            if cloned_bg is not None:
                svg.append(cloned_bg)

    return ET.tostring(svg, encoding="utf-8")


def _background_cache_key(svg_bytes: bytes, out_w: float, out_h: float) -> str:
    payload = f"{int(round(out_w))}x{int(round(out_h))}|".encode("ascii") + svg_bytes
    return hashlib.sha1(payload).hexdigest()


def _png_has_visible_pixels(png_bytes: bytes) -> bool:
    try:
        with Image.open(io.BytesIO(png_bytes)) as img:
            rgba = img.convert("RGBA")
            alpha = rgba.getchannel("A")
            extrema = alpha.getextrema()
            return bool(extrema and extrema[1] > 0)
    except Exception:
        return True


def find_background_base_rect(bg_group: ET.Element, vb_w: float, vb_h: float) -> Optional[ET.Element]:
    for child in bg_group:
        if local_name(child.tag) != "rect":
            continue
        if child.get("transform"):
            continue

        x = parse_length(child.get("x"), reference=vb_w, default=0.0)
        y = parse_length(child.get("y"), reference=vb_h, default=0.0)
        w = parse_length(child.get("width"), reference=vb_w, default=0.0)
        h = parse_length(child.get("height"), reference=vb_h, default=0.0)
        if x > 1e-3 or y > 1e-3 or w < vb_w - 1e-3 or h < vb_h - 1e-3:
            continue

        rect_style = parse_style_attr(child.get("style"))
        fill_text = child.get("fill") or rect_style.get("fill") or DEFAULT_STYLE["fill"]
        stroke_text = child.get("stroke") or rect_style.get("stroke") or DEFAULT_STYLE["stroke"]
        if fill_text.strip().lower().startswith("url("):
            continue
        if parse_paint_value(fill_text, {"color": rect_style.get("color", DEFAULT_STYLE["color"])}) is None:
            continue
        if stroke_text.strip().lower() not in {"", "none"}:
            continue
        if child.get("filter") or child.get("mask") or child.get("clip-path"):
            continue
        return child
    return None


def render_background_base_rect(bg_group: ET.Element, base_rect: ET.Element, ctx: RenderContext) -> None:
    group_style = merge_style(DEFAULT_STYLE, bg_group, ctx.style_rules)
    own_group_style = merge_style({}, bg_group, ctx.style_rules)
    group_opacity = parse_opacity_value(own_group_style.get("opacity"), default=1.0)
    group_transform = parse_transform(bg_group.get("transform"))
    render_element(base_rect, ctx, group_style, group_transform, group_opacity)


def render_background_group_image(
    root: ET.Element,
    bg_group: ET.Element,
    vb_w: float,
    vb_h: float,
    out_w: float,
    out_h: float,
    bg_cache: Optional[BackgroundImageCache] = None,
    omitted_ids: Optional[set[int]] = None,
) -> Optional[io.BytesIO]:
    svg_bytes = _build_background_svg_bytes(root, bg_group, vb_w, vb_h, omitted_ids=omitted_ids)
    cache_key = _background_cache_key(svg_bytes, out_w, out_h)

    if bg_cache is not None:
        cached_png = bg_cache.png_by_key.get(cache_key)
        if cached_png is not None:
            bg_cache.hits += 1
            return io.BytesIO(cached_png)

    png_bytes = _render_svg_bytes_to_png(svg_bytes, out_w, out_h)
    if not png_bytes or not _png_has_visible_pixels(png_bytes):
        return None

    if bg_cache is not None:
        bg_cache.png_by_key[cache_key] = png_bytes
        bg_cache.renders += 1

    return io.BytesIO(png_bytes)


def render_svg_to_slide(
    svg_path: Path,
    slide,
    target_width: float,
    target_height: float,
    use_bg_image: bool = True,
    bg_cache: Optional[BackgroundImageCache] = None,
) -> Tuple[float, float]:
    tree = ET.parse(svg_path)
    root = tree.getroot()

    vb_min_x, vb_min_y, vb_w, vb_h = parse_viewbox(root)
    gradients = parse_gradients(root, vb_w, vb_h)
    filters = parse_filters(root)
    style_rules = parse_css_rules(root)
    scale_x, scale_y, offset_x, offset_y, stroke_scale = compute_canvas_mapping(
        vb_w,
        vb_h,
        target_width,
        target_height,
    )

    ctx = RenderContext(
        slide=slide,
        vb_min_x=vb_min_x,
        vb_min_y=vb_min_y,
        vb_width=vb_w,
        vb_height=vb_h,
        target_width=target_width,
        target_height=target_height,
        scale_x=scale_x,
        scale_y=scale_y,
        offset_x=offset_x,
        offset_y=offset_y,
        stroke_scale=stroke_scale,
        svg_dir=svg_path.parent,
        gradients=gradients,
        filters=filters,
        style_rules=style_rules,
        skip_element_ids=set(),
    )

    if use_bg_image:
        bg_group = find_background_group(root)
        if bg_group is not None:
            omitted_ids: set[int] = set()
            base_rect = find_background_base_rect(bg_group, vb_w, vb_h)
            if base_rect is not None:
                render_background_base_rect(bg_group, base_rect, ctx)
                ctx.skip_element_ids.add(id(base_rect))
                omitted_ids.add(id(base_rect))

            bg_stream = render_background_group_image(
                root,
                bg_group,
                vb_w,
                vb_h,
                vb_w * scale_x,
                vb_h * scale_y,
                bg_cache=bg_cache,
                omitted_ids=omitted_ids,
            )
            if bg_stream is not None:
                slide.shapes.add_picture(
                    bg_stream,
                    Emu(emu(offset_x)),
                    Emu(emu(offset_y)),
                    Emu(emu(vb_w * scale_x)),
                    Emu(emu(vb_h * scale_y)),
                )
                ctx.skip_element_ids.add(id(bg_group))

    render_element(root, ctx, DEFAULT_STYLE, matrix_identity(), 1.0)
    return vb_w, vb_h


def find_svg_inputs(input_path: Path) -> List[Path]:
    if input_path.is_file() and input_path.suffix.lower() == ".svg":
        return [input_path]
    if input_path.is_dir():
        return sorted(input_path.glob("*.svg"))
    return []


def _emit_log(log_fn: Optional[Callable[[str], None]], message: str) -> None:
    if log_fn is not None:
        log_fn(message)


def _raise_if_timed_out(start_time: float, timeout_seconds: Optional[float]) -> None:
    if timeout_seconds is None or timeout_seconds <= 0:
        return
    if time.monotonic() - start_time > timeout_seconds:
        raise TimeoutError(f"SVG to PPTX conversion timed out after {timeout_seconds:.0f} seconds")


def convert_svg_files_to_pptx(
    svg_files: List[Path],
    output_path: Path,
    use_bg_image: bool = True,
    strict: bool = False,
    log_fn: Optional[Callable[[str], None]] = print,
    timeout_seconds: Optional[float] = None,
) -> Dict[str, object]:
    if not svg_files:
        raise ValueError("no SVG files found")

    started_at = time.monotonic()
    prs = Presentation()
    root0 = ET.parse(svg_files[0]).getroot()
    _, _, first_w, first_h = parse_viewbox(root0)
    prs.slide_width = Emu(emu(first_w))
    prs.slide_height = Emu(emu(first_h))
    target_width = float(first_w)
    target_height = float(first_h)

    blank = prs.slide_layouts[6]

    effective_bg_image = use_bg_image
    if effective_bg_image and not (HAS_CAIROSVG or HAS_SKIA):
        _emit_log(log_fn, "Warning: no SVG raster backend available, background remains editable vectors.")
        effective_bg_image = False

    bg_cache = BackgroundImageCache() if effective_bg_image else None
    failed_slides: List[Dict[str, object]] = []
    for idx, svg_path in enumerate(svg_files, start=1):
        _raise_if_timed_out(started_at, timeout_seconds)
        slide = prs.slides.add_slide(blank)
        try:
            w, h = render_svg_to_slide(
                svg_path,
                slide,
                target_width=target_width,
                target_height=target_height,
                use_bg_image=effective_bg_image,
                bg_cache=bg_cache,
            )
            _emit_log(log_fn, f"[{idx}/{len(svg_files)}] OK {svg_path.name} ({int(w)}x{int(h)} -> {int(target_width)}x{int(target_height)})")
        except Exception as exc:
            _emit_log(log_fn, f"[{idx}/{len(svg_files)}] FAIL {svg_path.name}: {exc}")
            failed_slides.append(
                {
                    "index": idx - 1,
                    "name": svg_path.name,
                    "error": str(exc),
                }
            )
            if strict:
                raise RuntimeError(f"slide conversion failed: {svg_path.name}: {exc}") from exc

    _raise_if_timed_out(started_at, timeout_seconds)
    if bg_cache is not None:
        _emit_log(
            log_fn,
            (
                "Background cache: "
                f"{bg_cache.hits} hits, {bg_cache.renders} renders, "
                f"{len(bg_cache.png_by_key)} unique backgrounds."
            ),
        )

    _raise_if_timed_out(started_at, timeout_seconds)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)

    return {
        "output_path": str(output_path),
        "total": len(svg_files),
        "failed": failed_slides,
        "success": len(svg_files) - len(failed_slides),
        "background_cache": {
            "enabled": effective_bg_image,
            "hits": bg_cache.hits if bg_cache is not None else 0,
            "renders": bg_cache.renders if bg_cache is not None else 0,
            "unique": len(bg_cache.png_by_key) if bg_cache is not None else 0,
        },
    }


def main() -> int:
    parser = argparse.ArgumentParser(description="Convert SVG to editable PPTX by reconstructing elements")
    parser.add_argument("input", type=str, help="SVG file or directory")
    parser.add_argument("-o", "--output", type=str, default=None, help="Output .pptx path")
    parser.add_argument(
        "--no-bg-image",
        action="store_true",
        help="Do not rasterize background group into a picture",
    )
    args = parser.parse_args()

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"Error: input does not exist: {input_path}")
        return 1

    svg_files = find_svg_inputs(input_path)
    if not svg_files:
        print("Error: no SVG files found")
        return 1

    if args.output:
        output_path = Path(args.output)
    else:
        if input_path.is_file():
            output_path = input_path.with_suffix(".editable.pptx")
        else:
            output_path = input_path / "editable_output.pptx"

    result = convert_svg_files_to_pptx(
        svg_files=svg_files,
        output_path=output_path,
        use_bg_image=not args.no_bg_image,
        strict=False,
        log_fn=print,
    )
    print(f"Saved: {result['output_path']}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
